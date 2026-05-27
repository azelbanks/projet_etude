#!/usr/bin/env python3
"""
Générateur PowerPoint 13 slides — Thumalien V9
Selon specs A_slides_keynote.md

Palette :  fond #0e1117, primaire #00D4FF (cyan), succès #00E676 (vert),
           danger #FF1744 (rouge), accent #FFD600 (jaune), texte #E8E8E8
Police :   Inter / SF Pro Display — titres 48-64pt, corps 28-32pt
Règle :    1 message / slide, max 3 puces, ≤ 30 mots simultanés
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ╔═══════════════════════════════════════════════════════════════════╗
# ║                         PALETTE                                  ║
# ╚═══════════════════════════════════════════════════════════════════╝

BG        = RGBColor(0x0E, 0x11, 0x17)   # Noir bleuté dashboard
BG_LIGHT  = RGBColor(0x14, 0x19, 0x24)   # Fond cartes légèrement plus clair
CARD_BG   = RGBColor(0x1A, 0x1F, 0x2E)   # Fond cartes
CYAN      = RGBColor(0x00, 0xD4, 0xFF)   # Primaire
CYAN_DIM  = RGBColor(0x00, 0x8A, 0xAA)   # Cyan atténué
GREEN     = RGBColor(0x00, 0xE6, 0x76)   # Succès
RED       = RGBColor(0xFF, 0x17, 0x44)   # Danger
YELLOW    = RGBColor(0xFF, 0xD6, 0x00)   # Accent
WHITE     = RGBColor(0xE8, 0xE8, 0xE8)   # Texte principal
GRAY      = RGBColor(0xB0, 0xB0, 0xB0)   # Texte secondaire
DARK_GRAY = RGBColor(0x60, 0x60, 0x60)   # Footer / discret
BORDER    = RGBColor(0x2A, 0x2F, 0x3E)   # Bordures subtiles
BLACK     = RGBColor(0x00, 0x00, 0x00)

FONT_MAIN = "Inter"

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)


# ╔═══════════════════════════════════════════════════════════════════╗
# ║                     HELPERS DE MISE EN PAGE                      ║
# ╚═══════════════════════════════════════════════════════════════════╝

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_notes(slide, text):
    """Ajoute des notes de présentation (speaker notes) à la slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_footer(slide, text="Niamato Consulting  ·  Thumalien V9  ·  M1 BDIA 2026"):
    tb = slide.shapes.add_textbox(Inches(9.0), Inches(7.0), Inches(4.0), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_MAIN
    p.alignment = PP_ALIGN.RIGHT


def add_slide_number(slide, num, total=17):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(1.2), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_MAIN
    p.alignment = PP_ALIGN.LEFT


def _fmt_run(paragraph, text, size, color, bold=False, italic=False, font=FONT_MAIN):
    """Ajoute un run formaté à un paragraphe existant."""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return run


def add_title(slide, text, color=CYAN, size=44, left=0.8, top=0.25, width=11.5):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = FONT_MAIN
    p.space_after = Pt(0)
    return tb


def add_subtitle(slide, text, left=0.8, top=0.95, size=22):
    return add_text(slide, text, left, top, 11, 0.4, size=size, color=GRAY, italic=False)


def add_text(slide, text, left, top, width, height, size=22, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, italic=False, font=FONT_MAIN,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Anchor vertical
    try:
        tf.paragraphs[0].alignment = align
    except Exception:
        pass
    # Gérer le multi-lignes
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = font
        p.alignment = align
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    return tb


def add_card(slide, left, top, width, height, fill_color=CARD_BG,
             border_color=CYAN, border_width=1.2, radius=0.15):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_width)
    # Ajuster le rayon de l'arrondi
    shape.adjustments[0] = radius
    return shape


def add_line_h(slide, left, top, width, color=CYAN, thickness=2):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Pt(thickness)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line_v(slide, left, top, height, color=CYAN, thickness=2):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Pt(thickness), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_icon_text(slide, icon, kw, desc, x, y, w, col, desc_color=GRAY):
    """Icône + mot-clé bold + description sur une carte horizontale."""
    add_card(slide, x, y, w, 0.95, CARD_BG, col, border_width=1.5)
    add_text(slide, icon, x + 0.25, y + 0.15, 0.6, 0.6, size=28, color=col,
             align=PP_ALIGN.CENTER)
    add_text(slide, kw, x + 1.0, y + 0.12, 3.0, 0.4,
             size=28, color=col, bold=True)
    add_text(slide, desc, x + 4.3, y + 0.18, w - 4.6, 0.55,
             size=20, color=desc_color)


def add_stat_card(slide, x, y, w, h, label, big, sub, col):
    """Panneau statistique : label en haut, gros chiffre, sous-texte."""
    add_card(slide, x, y, w, h, CARD_BG, col, border_width=1.5)
    add_text(slide, label, x + 0.15, y + 0.2, w - 0.3, 0.35,
             size=18, color=col, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, big, x + 0.1, y + 0.7, w - 0.2, 1.2,
             size=80, color=col, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, sub, x + 0.2, y + 2.2, w - 0.4, h - 2.5,
             size=17, color=GRAY, align=PP_ALIGN.CENTER)


# ╔═══════════════════════════════════════════════════════════════════╗
# ║                    CRÉATION DES 14 SLIDES                        ║
# ╚═══════════════════════════════════════════════════════════════════╝

def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    n = 0  # compteur de slides

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 1 — TITRE PROJET                                    00:30
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)

    # Filet décoratif haut
    add_line_h(slide, 2.5, 1.2, 8.3, CYAN_DIM, 1)

    add_text(slide, "THUMALIEN V9", 0.5, 1.5, 12.3, 1.3,
             size=76, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, "Détection de désinformation Bluesky — bilingue, explicable, frugale",
             1.0, 3.0, 11.3, 0.65, size=26, color=GRAY, align=PP_ALIGN.CENTER)

    add_line_h(slide, 3.5, 3.95, 6.3, CYAN, 2)

    add_text(slide, "NIAMATO CONSULTING",
             1.0, 4.2, 11.3, 0.45, size=22, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Azélie Bernard  ·  Sébastien Lazcanotegui",
             1.0, 4.7, 11.3, 0.55, size=26, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Master 1 Big Data & Intelligence Artificielle — 2025 / 2026",
             1.0, 5.3, 11.3, 0.4, size=18, color=GRAY, align=PP_ALIGN.CENTER)

    # Filet décoratif bas
    add_line_h(slide, 2.5, 6.3, 8.3, CYAN_DIM, 1)

    add_text(slide, "[QR code → repo GitHub]", 10.5, 5.8, 2.5, 0.4,
             size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 1 — TITRE (00:00 → 00:30)
Speaker : AZÉLIE (hook)
Timing : 30 secondes

Texte voix-off (Azélie) :
« En décembre 2025, on collecte 100 000 posts Bluesky.
On entraîne un modèle. On obtient un F1 de 0,99.
Et c'est précisément à ce moment-là qu'on aurait dû s'inquiéter.

Bonjour, je suis Azélie Bernard, lead technique sur Thumalien.
Avec Sébastien, en 18 minutes on va vous raconter
pourquoi ce 0,99 était un piège. »

Points clés :
- Ouverture contre-intuitive = capte l'attention
- Regarder la caméra, ton grave puis sourire
- BANDEAU OBS : Azélie Bernard — Lead Technique""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 2 — ÉQUIPE ET RÔLES                                 00:30
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "NIAMATO CONSULTING — L'ÉQUIPE")

    # Carte Azélie
    add_card(slide, 0.8, 1.5, 5.5, 4.8, CARD_BG, CYAN, 1.5)
    add_text(slide, "[PHOTO]", 1.3, 1.8, 1.8, 1.8,
             size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_card(slide, 1.3, 1.8, 1.8, 1.8, CARD_BG, CYAN_DIM, 1)
    add_text(slide, "Azélie Bernard", 3.5, 1.8, 2.5, 0.5,
             size=24, color=CYAN, bold=True)
    add_text(slide, "Lead Technique", 3.5, 2.35, 2.5, 0.4,
             size=18, color=YELLOW, bold=True)
    add_line_h(slide, 1.2, 3.9, 4.8, BORDER, 1)
    roles_a = [
        "Pipeline NLP V1 \u2192 V9",
        "Dashboard Streamlit 5 pages",
        "XAI : SHAP, Captum, faithfulness",
        "Fine-tuning CamemBERT & RoBERTa",
        "CI/CD, tests, mutation testing",
    ]
    for j, r in enumerate(roles_a):
        add_text(slide, f"\u2022  {r}", 1.2, 4.1 + j * 0.4, 4.8, 0.35,
                 size=15, color=WHITE)

    # Carte Sébastien
    add_card(slide, 7.0, 1.5, 5.5, 4.8, CARD_BG, GREEN, 1.5)
    add_text(slide, "[PHOTO]", 7.5, 1.8, 1.8, 1.8,
             size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_card(slide, 7.5, 1.8, 1.8, 1.8, CARD_BG, CYAN_DIM, 1)
    add_text(slide, "Sébastien Lazcanotegui", 9.7, 1.8, 2.6, 0.5,
             size=22, color=GREEN, bold=True)
    add_text(slide, "Validation & Qualité ML", 9.7, 2.35, 2.6, 0.4,
             size=16, color=YELLOW, bold=True)
    add_line_h(slide, 7.4, 3.9, 4.8, BORDER, 1)
    roles_s = [
        "Annotation gold set (473 posts, 2 annotateurs)",
        "Débiaisage Reuters (BODY_AGENCY_TERMS)",
        "GridSearch hyperparamètres (C, min_df, ngrams)",
        "Revue critique & challenge technique",
        "Conformité RGPD & AI Act",
    ]
    for j, r in enumerate(roles_s):
        add_text(slide, f"\u2022  {r}", 7.4, 4.1 + j * 0.4, 4.8, 0.35,
                 size=15, color=WHITE)

    add_text(slide, "Master 1 Big Data & IA — Sup de Vinci — 2025/2026",
             0.5, 6.6, 12.3, 0.35, size=16, color=GRAY, italic=True,
             align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 2 \u2014 ÉQUIPE ET RÔLES (00:30)
Timing : 30 secondes (Sébastien se présente)

Texte voix-off (Sébastien) :
\u00ab Bonjour, je suis Sébastien Lazcanotegui.
Mon rôle : la validation et la qualité du machine learning.
J'ai annoté les 473 posts du gold set avec Azélie,
piloté le débiaisage du biais Reuters,
et optimisé les hyperparamètres par GridSearch.
Ma collègue Azélie Bernard est la lead technique \u2014
elle a conçu et itéré l'ensemble du pipeline,
du collecteur Bluesky au dashboard en passant par l'XAI. \u00bb

IMPORTANT : insérer les photos avant export final.""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 3 — PROBLÉMATIQUE CHIFFRÉE                          01:00
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "LA DÉSINFORMATION SUR BLUESKY")

    stats = [
        ("35 M",      "utilisateurs Bluesky",    CYAN),
        ("60 000+",   "posts publics FR/EN par jour", CYAN),
        ("0",         "équipe modération centralisée", RED),
        ("Angle mort","régulation européenne",    YELLOW),
    ]
    positions = [(0.8, 1.6), (6.8, 1.6), (0.8, 4.0), (6.8, 4.0)]
    for (val, label, col), (x, y) in zip(stats, positions):
        add_card(slide, x, y, 5.7, 2.0, CARD_BG, col, border_width=1.5)
        add_text(slide, val, x + 0.3, y + 0.2, 5.1, 1.0,
                 size=64, color=col, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.3, y + 1.25, 5.1, 0.55,
                 size=20, color=GRAY, align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 2 — PROBLÉMATIQUE CHIFFRÉE (00:30 → 01:00)
Timing : 30 secondes
Animation : 4 cellules en cascade haut-gauche → bas-droite, intervalle 0.3s

Texte voix-off :
« Bluesky compte 35 millions d'utilisateurs. Plus de 60 000 posts publics
en français et anglais chaque jour. Et pourtant : zéro équipe de modération
centralisée, et un angle mort complet côté régulation européenne.
C'est cet espace non couvert que Thumalien vient adresser. »

Point d'insistance : marquer une pause sur "zéro équipe" et "angle mort".""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 3 — LES 4 EXIGENCES DU CAHIER DES CHARGES           02:00
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "LE CAHIER DES CHARGES")

    exigences = [
        ("\U0001F50D", "TRANSPARENCE",  "Chaque score doit être explicable et auditable",    CYAN),
        ("\U0001F310", "BILINGUE",      "Français + anglais, performance équivalente",        GREEN),
        ("\U0001F331", "FRUGALITÉ",     "< 5 ms par texte, empreinte CO₂ mesurée",           YELLOW),
        ("\u2696\uFE0F","CONFORMITÉ",   "RGPD art. 22 + AI Act art. 13/14",                  RED),
    ]
    for i, (icon, kw, desc, col) in enumerate(exigences):
        y = 1.5 + i * 1.25
        add_icon_text(slide, icon, kw, desc, 0.8, y, 11.7, col)

    add_text(slide,
             "Quatre exigences non-négociables pour passer du prototype au système déployable",
             0.8, 6.6, 11.7, 0.35, size=16, color=CYAN, italic=True,
             align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 3 — CAHIER DES CHARGES (01:00 → 02:00)
Timing : 60 secondes
Animation : cartes apparaissent une par une, 1 seconde d'intervalle

Texte voix-off :
« Notre cahier des charges impose quatre exigences non-négociables.
Premièrement, la transparence : chaque prédiction doit être explicable.
Deuxièmement, le bilinguisme : français et anglais à performance identique.
Troisièmement, la frugalité : moins de 5 millisecondes par texte
avec une empreinte carbone mesurée.
Enfin, la conformité réglementaire : RGPD article 22 et AI Act articles 13-14. »

Point clé : ces 4 exigences structurent TOUTE la présentation.""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 4 — LA FAUSSE VICTOIRE (F1=0.99)                    03:00
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "F1 = 0.99 — LE PIÈGE", RED, size=48)

    # ── Panneau gauche : métriques parfaites ──
    add_card(slide, 0.6, 1.4, 5.7, 4.8, CARD_BG, GREEN, 1.5)
    add_text(slide, "CROSS-VALIDATION 5-FOLD", 1.0, 1.6, 4.8, 0.4,
             size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_line_h(slide, 1.2, 2.15, 4.5, BORDER, 1)

    metrics = [
        ("F1 macro",  "0.987"),
        ("Précision", "0.992"),
        ("Recall",    "0.984"),
        ("Accuracy",  "0.989"),
    ]
    for j, (name, val) in enumerate(metrics):
        y = 2.4 + j * 0.6
        add_text(slide, name, 1.2, y, 2.8, 0.4, size=22, color=GRAY)
        add_text(slide, val, 4.2, y, 1.8, 0.4, size=24, color=GREEN,
                 bold=True, align=PP_ALIGN.RIGHT)

    add_text(slide, "✅  Tous les voyants au vert", 1.0, 5.0, 4.8, 0.4,
             size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # ── Séparateur vertical rouge ──
    add_line_v(slide, 6.6, 1.4, 4.8, RED, 3)

    # ── Panneau droit : top mots révélateurs ──
    add_card(slide, 7.0, 1.4, 5.7, 4.8, CARD_BG, RED, 1.5)
    add_text(slide, "TOP COEFFICIENTS LOGREG (XAI)", 7.3, 1.6, 5.0, 0.4,
             size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_line_h(slide, 7.4, 2.15, 4.9, BORDER, 1)

    words = [
        ("reuters",     "+0.84"),
        ("afp",         "+0.71"),
        ("ap_news",     "+0.69"),
        ("associated",  "+0.63"),
        ("press",       "+0.58"),
    ]
    for j, (w, s) in enumerate(words):
        y = 2.4 + j * 0.5
        add_text(slide, w, 7.6, y, 3.0, 0.4, size=22, color=WHITE,
                 font="Courier New")
        add_text(slide, s, 10.8, y, 1.5, 0.4, size=22, color=RED,
                 bold=True, align=PP_ALIGN.RIGHT, font="Courier New")

    add_text(slide, "⚠️  Le modèle apprend le STYLE\n     des agences de presse, pas les faits",
             7.3, 5.0, 5.0, 0.7, size=16, color=YELLOW, bold=True)

    # ── Footer dramatique ──
    add_text(slide,
             "Sans XAI, on aurait livré un modèle qui collapse en production",
             0.6, 6.5, 12.1, 0.35, size=18, color=RED, italic=True,
             align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 4 — F1=0.99 LE PIÈGE (02:00 → 03:00)
Timing : 60 secondes
Animation CRUCIALE : le panneau gauche apparaît immédiatement.
Le panneau droit arrive 3 secondes après — c'est la RÉVÉLATION.

Texte voix-off :
« Nos premières métriques étaient extraordinaires : F1 de 0.987, précision 0.99.
Tous les voyants étaient au vert. [PAUSE 2s]
Mais quand on a ouvert la boîte noire avec l'explicabilité…
les coefficients les plus importants étaient : "reuters", "afp", "associated press".
Le modèle n'apprenait pas à détecter la désinformation.
Il apprenait à reconnaître le style des agences de presse.
Sans XAI, on aurait livré un modèle inutile en production. »

C'est LE moment clé de la présentation — le ton doit devenir grave.""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 5 — ARCHITECTURE C4 NIVEAU 2                        05:00
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "ARCHITECTURE — 8 CONTENEURS")

    boxes = [
        ("Bluesky\nAT Proto",          0.5,  1.4, 3.3, 1.4, CYAN),
        ("Collector\n(Python)",         4.5,  1.4, 3.3, 1.4, CYAN),
        ("MongoDB\n(storage)",          8.5,  1.4, 3.3, 1.4, CYAN),
        ("Detector\nV5 + V6 + CamBERT",8.5,  3.4, 3.3, 1.6, GREEN),
        ("Pipeline Cascade\nStage 1+2", 4.5,  3.4, 3.3, 1.6, GREEN),
        ("Streamlit\nDashboard",        0.5,  3.4, 3.3, 1.6, YELLOW),
        ("XAI Engine\nSHAP + Captum",   0.5,  5.6, 3.3, 1.2, CYAN_DIM),
        ("Monitoring\nCodeCarbon",      4.5,  5.6, 3.3, 1.2, CYAN_DIM),
    ]
    for label, x, y, w, h, col in boxes:
        card = add_card(slide, x, y, w, h, CARD_BG, col, 1.2)
        tf = card.text_frame
        tf.word_wrap = True
        lines = label.split('\n')
        p = tf.paragraphs[0]
        p.text = lines[0]
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT_MAIN
        p.alignment = PP_ALIGN.CENTER
        for extra in lines[1:]:
            p2 = tf.add_paragraph()
            p2.text = extra
            p2.font.size = Pt(14)
            p2.font.color.rgb = GRAY
            p2.font.name = FONT_MAIN
            p2.alignment = PP_ALIGN.CENTER

    # Flèches horizontales (simplifiées)
    for (x1, x2, y) in [(3.8, 4.5, 2.1), (7.8, 8.5, 2.1),
                          (7.8, 8.5, 4.2), (3.8, 4.5, 4.2),
                          (3.8, 4.5, 6.2)]:
        add_text(slide, "→", x1, y, x2 - x1, 0.3, size=24, color=CYAN,
                 align=PP_ALIGN.CENTER)

    # Flèches verticales
    for (x, y) in [(9.8, 2.8), (6.0, 2.8), (1.8, 5.0)]:
        add_text(slide, "↓", x, y, 0.5, 0.4, size=22, color=CYAN,
                 align=PP_ALIGN.CENTER)

    add_text(slide, "Données → IA → Décision  |  C4 Model  |  Docker Compose  |  FastAPI",
             7.0, 7.0, 6.0, 0.3,
             size=13, color=DARK_GRAY, align=PP_ALIGN.RIGHT, italic=True)
    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 5 — ARCHITECTURE C4 (03:00 → 05:00)
Timing : 2 minutes
Animation : boîtes apparaissant de gauche à droite, haut en bas

Texte voix-off :
« Voici notre architecture en 8 conteneurs, selon le modèle C4 niveau 2.
Le flux commence par l'API Bluesky AT Proto, collecte les posts via un
collecteur Python, et les stocke dans MongoDB.
Le pipeline cascade en deux étapes filtre d'abord les opinions
puis analyse les faits avec nos trois modèles : V5 TF-IDF, V6 style,
et CamemBERT.
Le dashboard Streamlit expose les résultats avec l'explicabilité
SHAP et Captum. CodeCarbon mesure l'empreinte carbone. »

Pourquoi C4 : standard industriel → montre maturité architecturale.""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 6 — PIPELINE CASCADE INFÉRENCE                      06:00
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "PIPELINE CASCADE — INFÉRENCE")

    actors = ["User", "Stage 1", "V5", "V6", "CamBERT", "V8 Méta", "XAI"]
    actor_w = 1.5
    start_x = 0.5
    for i, actor in enumerate(actors):
        x = start_x + i * 1.75
        # Boîte acteur
        add_card(slide, x, 1.3, actor_w, 0.55, CARD_BG, CYAN, 1)
        add_text(slide, actor, x, 1.33, actor_w, 0.5,
                 size=14, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
        # Ligne de vie
        add_line_v(slide, x + actor_w / 2, 1.85, 4.5, BORDER, 1)

    # Messages séquence
    msgs = [
        (0, 1, "texte",       2.1),
        (1, 1, "opinion ?",   2.6),
        (1, 2, "si factuel",  3.1),
        (2, 3, "score TF-IDF",3.5),
        (3, 4, "score style", 3.9),
        (4, 5, "embedding",   4.3),
        (5, 6, "décomposition", 4.7),
        (5, 1, "score V8",    5.2),
        (1, 0, "score + XAI", 5.6),
    ]
    for src, dst, label, y in msgs:
        x1 = start_x + src * 1.75 + actor_w / 2
        x2 = start_x + dst * 1.75 + actor_w / 2
        left = min(x1, x2)
        w = abs(x2 - x1) if x1 != x2 else 0.4
        # Ligne
        add_line_h(slide, left, y + 0.1, w, CYAN, 2)
        # Étiquette
        add_text(slide, label, left, y - 0.15, w, 0.25,
                 size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # Badge latence
    add_card(slide, 9.5, 6.0, 3.0, 0.7, CARD_BG, YELLOW, 1.5)
    add_text(slide, "⏱  1.5 ms total", 9.5, 6.05, 3.0, 0.6,
             size=20, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)

    # Badge AI Act
    add_text(slide, "⚖️ opinion ≠ fake news", 0.5, 6.2, 3.5, 0.35,
             size=14, color=GRAY, italic=True)

    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 6 — PIPELINE CASCADE INFÉRENCE (05:00 → 06:00)
Timing : 60 secondes
Animation : flèches défilant une par une, 0.3s d'intervalle

Texte voix-off :
« Le pipeline fonctionne en cascade. D'abord, le Stage 1 détecte les opinions
— conformément à l'AI Act, une opinion n'est pas de la désinformation.
Si le texte est factuel, il passe par V5 pour le scoring TF-IDF,
puis V6 pour l'analyse stylistique, puis CamemBERT pour l'embedding sémantique.
Le méta-modèle V8 combine les trois scores, et le module XAI fournit
l'explication complète.
Le tout en 1.5 milliseconde par texte. »""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 7 — FAITHFULNESS : POURQUOI ET COMMENT              10:30
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "LA FIDÉLITÉ DES EXPLICATIONS")
    add_subtitle(slide, "Au-delà de SHAP : la mesurer")

    # Question
    add_card(slide, 0.8, 1.7, 11.7, 1.5, CARD_BG, CYAN, 1.5)
    add_text(slide,
             "❓ Comment savez-vous que vos explications\n"
             "    reflètent vraiment le comportement du modèle ?",
             1.3, 1.85, 10.7, 1.2, size=26, color=CYAN, align=PP_ALIGN.CENTER)

    # Flèche
    add_text(slide, "⬇", 6.2, 3.3, 1, 0.5, size=36, color=CYAN,
             align=PP_ALIGN.CENTER)

    # Méthode
    add_card(slide, 0.8, 3.9, 11.7, 2.9, CARD_BG, CYAN, 1.5)
    add_text(slide, "PROTOCOLE ERASER  (DeYoung et al., ACL 2020)",
             1.3, 4.1, 10.5, 0.4, size=20, color=WHITE, bold=True)
    add_line_h(slide, 1.3, 4.6, 10.5, BORDER, 1)

    steps = [
        "1.  Identifier les top-k features via SHAP",
        "2.  Masquer ces features (valeur = 0)",
        "3.  Mesurer la chute de P(suspect)",
        "4.  Comparer à un masquage aléatoire",
    ]
    for j, step in enumerate(steps):
        add_text(slide, step, 1.8, 4.8 + j * 0.4, 9.5, 0.35,
                 size=19, color=GRAY)

    add_text(slide,
             "Si fidèle :  chute_attribution  ≫  chute_random",
             1.3, 6.4, 10.5, 0.35, size=22, color=YELLOW, bold=True,
             align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 7 — FAITHFULNESS MÉTHODE (06:00 → 10:30)
Timing : cette slide sert de support pendant la démo live ou l'explication longue

Texte voix-off :
« Avoir SHAP ne suffit pas. Il faut prouver que les explications sont fidèles.
Notre question : si on supprime les features que SHAP identifie comme
les plus importantes, est-ce que la prédiction change vraiment ?
On utilise le protocole ERASER de DeYoung et al., publié à ACL 2020.
En 4 étapes : identifier, masquer, mesurer, comparer.
Si notre attribution est fidèle, la chute de probabilité doit être
nettement supérieure à un masquage aléatoire. »

Référence académique : important pour la crédibilité.""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 8 — COURBE AOPC                                     11:30
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "RÉSULTAT FAITHFULNESS — UPLIFT +21%", GREEN, size=40)

    # Image AOPC
    aopc_path = os.path.join(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
        "docs", "figures", "xai", "faithfulness_aopc_curve.png"
    )
    if os.path.exists(aopc_path):
        slide.shapes.add_picture(aopc_path, Inches(0.5), Inches(1.4),
                                  Inches(8.2), Inches(5.3))
    else:
        add_card(slide, 0.5, 1.4, 8.2, 5.3, CARD_BG, CYAN, 1)
        add_text(slide, "[Insérer faithfulness_aopc_curve.png]",
                 1.5, 3.8, 6, 0.5, size=22, color=GRAY, align=PP_ALIGN.CENTER)

    # Panneau stats à droite
    add_card(slide, 9.0, 1.4, 4.0, 5.3, CARD_BG, GREEN, 1.5)
    add_text(slide, "MÉTRIQUES", 9.2, 1.6, 3.6, 0.35,
             size=16, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_line_h(slide, 9.3, 2.05, 3.4, BORDER, 1)

    kv = [
        ("AOPC attrib.",  "0.253",  WHITE),
        ("AOPC random",   "0.045",  GRAY),
    ]
    y = 2.2
    for lab, val, vcol in kv:
        add_text(slide, lab, 9.3, y, 2.2, 0.3, size=15, color=GRAY)
        add_text(slide, val, 11.6, y, 1.2, 0.3,
                 size=17, color=vcol, bold=True, align=PP_ALIGN.RIGHT)
        y += 0.35

    add_line_h(slide, 9.3, y + 0.05, 3.4, BORDER, 1)
    y += 0.2

    kv2 = [
        ("Uplift",           "+0.21",  GREEN),
        ("Ratio",            "5.6×",   GREEN),
    ]
    for lab, val, vcol in kv2:
        add_text(slide, lab, 9.3, y, 2.2, 0.3, size=15, color=GRAY)
        add_text(slide, val, 11.6, y, 1.2, 0.3,
                 size=18, color=vcol, bold=True, align=PP_ALIGN.RIGHT)
        y += 0.38

    add_line_h(slide, 9.3, y + 0.05, 3.4, BORDER, 1)
    y += 0.2

    kv3 = [
        ("Compr.@5", "0.232"),
        ("Suff.@5",  "0.058"),
    ]
    for lab, val in kv3:
        add_text(slide, lab, 9.3, y, 2.2, 0.3, size=15, color=GRAY)
        add_text(slide, val, 11.6, y, 1.2, 0.3,
                 size=16, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
        y += 0.35

    add_text(slide, "✅ Cible > +0.10\n✅ Atteinte (× 2)", 9.3, y + 0.15, 3.4, 0.6,
             size=16, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 8 — RÉSULTAT FAITHFULNESS (10:30 → 11:30)
Timing : 60 secondes

Texte voix-off :
« Voici le résultat concret. La courbe bleue montre la chute de probabilité
quand on masque les features dans l'ordre d'importance SHAP.
La courbe orange, c'est le masquage aléatoire.
L'AOPC de notre attribution est de 0.253 contre 0.045 pour le random.
C'est un uplift de 21 points, un ratio de 5.6 fois.
Notre cible était un uplift supérieur à 10 points — on l'atteint
avec un facteur 2. Nos explications sont fidèles au modèle. »""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 9 — QUALITÉ INDUSTRIELLE                             12:30
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "QUALITÉ INDUSTRIELLE")
    add_subtitle(slide, "501 tests  ·  80 % coverage  ·  80,3 % mutation kill rate")

    add_stat_card(slide, 0.6, 1.6, 3.8, 5.0,
                  "PYTEST", "501", "tests passing\n\n✅ 0 failure\n✅ 0 skip", CYAN)
    add_stat_card(slide, 4.8, 1.6, 3.8, 5.0,
                  "COVERAGE", "80 %",
                  "line coverage\n77,9 % branch cov\n\nQuality gate CI\n--fail-under=75", GREEN)
    add_stat_card(slide, 9.0, 1.6, 3.8, 5.0,
                  "MUTMUT", "80,3 %",
                  "kill rate\n143 / 178 mutants\n\n> benchmark Google\n(60–75 %)", YELLOW)

    add_text(slide,
             "De la quantité à la qualité — chaque ligne de production est validée",
             0.6, 6.9, 12.1, 0.3, size=16, color=GRAY, italic=True,
             align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 9 — QUALITÉ INDUSTRIELLE (11:30 → 12:30)
Timing : 60 secondes

Texte voix-off :
« 501 tests unitaires, tous passants. 80 % de couverture de code.
Mais la couverture seule ne prouve rien — c'est pour ça qu'on a lancé
le mutation testing avec mutmut.
80,3 % des mutants sont tués. Pour contexte, le benchmark industriel
chez Google est entre 60 et 75 %. On est au-dessus.
Et on a mis en place un quality gate dans la CI : aucun merge
n'est accepté si la couverture descend sous 75 %. »

Point d'argumentation : mutation testing = preuve que les tests
détectent VRAIMENT les bugs, pas juste qu'ils traversent les lignes.""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 10 — CONFORMITÉ AI ACT & RGPD                       14:00
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "CONFORMITÉ RÉGLEMENTAIRE")

    # ── AI Act ──
    add_text(slide, "AI ACT  (UE 2024/1689)", 0.8, 1.4, 6, 0.4,
             size=22, color=CYAN, bold=True)
    add_text(slide, "pleinement applicable août 2026", 7.5, 1.45, 5, 0.35,
             size=16, color=GRAY, italic=True)

    add_line_h(slide, 0.8, 1.9, 11.7, BORDER, 1)

    ai_items = [
        ("Art. 13 (transparence)", "→  Model Card MC-THUM-2026-001"),
        ("Art. 14 (supervision)",  "→  Décomposition β·x dans le dashboard"),
        ("Risque limité classifié","→  Doc 02 § 4.1"),
    ]
    for j, (art, proof) in enumerate(ai_items):
        y = 2.1 + j * 0.55
        add_text(slide, "✅", 0.8, y, 0.4, 0.4, size=18, color=GREEN)
        add_text(slide, art, 1.4, y, 4.5, 0.4, size=19, color=WHITE)
        add_text(slide, proof, 6.5, y, 6, 0.4, size=17, color=GRAY,
                 font="Courier New")

    # ── Séparateur ──
    add_line_h(slide, 0.8, 3.9, 11.7, CYAN_DIM, 1)

    # ── RGPD ──
    add_text(slide, "RGPD  (UE 2016/679)", 0.8, 4.15, 6, 0.4,
             size=22, color=CYAN, bold=True)

    add_line_h(slide, 0.8, 4.65, 11.7, BORDER, 1)

    rgpd_items = [
        ("Art. 22 (décision auto)", "→  Droit à l'explication SHAP + Captum"),
        ("Art. 35 (AIPD)",          "→  Document RGPD-THUM-2026-001"),
        ("Base légale art. 6.1.f",  "→  Posts publics, intérêt légitime"),
    ]
    for j, (art, proof) in enumerate(rgpd_items):
        y = 4.85 + j * 0.55
        add_text(slide, "✅", 0.8, y, 0.4, 0.4, size=18, color=GREEN)
        add_text(slide, art, 1.4, y, 4.5, 0.4, size=19, color=WHITE)
        add_text(slide, proof, 6.5, y, 6, 0.4, size=17, color=GRAY,
                 font="Courier New")

    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 10 — CONFORMITÉ RÉGLEMENTAIRE (12:30 → 14:00)
Timing : 90 secondes

Texte voix-off :
« Thumalien est conforme aux deux cadres réglementaires européens.
Pour l'AI Act, applicable en août 2026 : notre Model Card satisfait l'article 13
sur la transparence. La décomposition β·x dans le dashboard répond à l'article 14
sur la supervision humaine. Et notre système est classé risque limité.
Côté RGPD : l'article 22 sur les décisions automatisées est couvert par
nos explications SHAP et Captum — chaque utilisateur peut comprendre
pourquoi un texte est flaggé. L'AIPD est documentée, et notre base légale
est l'intérêt légitime sur des posts publics. »""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 11 — GREEN IT                                       14:30
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "GREEN IT — ≈ 6,88 g CO₂ TOTAL", GREEN)

    # ── Panneau gauche : barres de répartition ──
    add_card(slide, 0.5, 1.4, 6.0, 5.6, CARD_BG, GREEN, 1.5)
    add_text(slide, "RÉPARTITION EMPREINTE CO₂", 0.8, 1.6, 5.4, 0.35,
             size=17, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_line_h(slide, 1.0, 2.05, 5.0, BORDER, 1)

    bars = [
        ("RoBERTa EN (V1+V2)",       50, "3,47 g"),
        ("V5 LogReg (3 runs)",        32, "2,18 g"),
        ("CamemBERT fine-tune",        7, "0,48 g"),
        ("V6 GradientBoost*",          5, "~0,36 g"),
        ("Pipeline XAI*",              3, "~0,23 g"),
        ("Inférence cumulée*",         2, "~0,15 g"),
    ]
    bar_colors = [CYAN, CYAN_DIM, YELLOW, GREEN, GRAY, GRAY]
    for j, ((label, pct, grams), bcol) in enumerate(zip(bars, bar_colors)):
        y = 2.3 + j * 0.7
        # Label
        add_text(slide, label, 1.0, y - 0.05, 3.0, 0.3, size=14, color=WHITE)
        # Barre
        bar_w = pct / 100 * 4.5
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(y + 0.25),
            Inches(bar_w), Inches(0.28)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bcol
        shape.line.fill.background()
        shape.adjustments[0] = 0.5  # pill shape
        # Valeur
        add_text(slide, f"{pct} %  ({grams})", 1.1 + bar_w, y + 0.22, 2.0, 0.3,
                 size=13, color=GRAY)

    # ── Panneau droit : décision archi ──
    add_card(slide, 7.0, 1.4, 5.5, 5.2, CARD_BG, CYAN, 1.5)
    add_text(slide, "DÉCISION ARCHITECTURALE", 7.2, 1.6, 5.1, 0.35,
             size=17, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_line_h(slide, 7.3, 2.05, 4.9, BORDER, 1)

    sections = [
        ("Production temps réel",
         "→  V5 seul (1,5 ms, 0,6 g CO₂/jour)",
         GREEN),
        ("Analyse offline batch",
         "→  V8 avec CamemBERT\n     (0,48 g CO₂ entraînement)",
         CYAN),
        ("ROI Green IT",
         "→  Frugalité prod + puissance analyse\n→  Documenté Model Card § 8",
         YELLOW),
    ]
    y = 2.3
    for title, desc, col in sections:
        add_text(slide, title, 7.4, y, 4.8, 0.3,
                 size=18, color=col, bold=True)
        add_text(slide, desc, 7.6, y + 0.35, 4.6, 0.7,
                 size=15, color=GRAY)
        y += 1.2

    add_text(slide, "* estimé via taux CodeCarbon (7,55×10⁻⁷ kg CO₂/s, Apple M4 Pro)",
             0.5, 6.55, 6.0, 0.25, size=11, color=DARK_GRAY, italic=True)
    add_text(slide, "≈ 30 mètres en voiture — empreinte quasi-nulle",
             0.5, 6.85, 12.3, 0.3, size=15, color=DARK_GRAY, italic=True,
             align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 11 — GREEN IT (14:00 → 14:30)
Timing : 30 secondes

Texte voix-off :
« Notre empreinte carbone totale est d'environ 6,9 grammes de CO2.
6,14 grammes mesurés par CodeCarbon sur 6 entraînements,
plus environ 0,7 gramme estimé pour V6, le pipeline XAI et l'inférence.
Les deux fine-tunings RoBERTa représentent la moitié du bilan.
Les trois runs LogReg un tiers, et le reste se répartit entre
CamemBERT, V6 GradientBoosting, l'explicabilité et l'inférence.
En production, on utilise V5 seul : 1,5 milliseconde, 0,6 gramme par jour.
CamemBERT est réservé à l'analyse offline batch.
C'est un arbitrage documenté dans la Model Card. »""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 13 — MÉTHODOLOGIE ET ORGANISATION                   15:00
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "MÉTHODOLOGIE & ORGANISATION")
    add_subtitle(slide, "CRISP-DM adapté  \u00b7  9 versions  \u00b7  28 notebooks  \u00b7  16 work packages")

    # Cycle CRISP-DM
    add_card(slide, 0.5, 1.5, 5.8, 5.0, CARD_BG, CYAN, 1.5)
    add_text(slide, "CYCLE CRISP-DM ADAPTÉ ML", 0.8, 1.7, 5.2, 0.35,
             size=17, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_line_h(slide, 0.9, 2.15, 5.0, BORDER, 1)

    crisp_steps = [
        ("1.", "Comprendre", "Cahier des charges, 4 exigences", CYAN),
        ("2.", "Explorer", "245K posts, audit qualité données", CYAN),
        ("3.", "Préparer", "Débiaisage Reuters, données synth.", GREEN),
        ("4.", "Modéliser", "V1\u2192V9, LogReg + CamemBERT + ensemble", GREEN),
        ("5.", "Évaluer", "Gold set 473 posts, faithfulness AOPC", YELLOW),
        ("6.", "Déployer", "Docker Compose, CI/CD, dashboard", YELLOW),
        ("\u21bb", "Itérer", "9 versions en 6 mois", RED),
    ]
    for j, (num, step, desc, col) in enumerate(crisp_steps):
        y = 2.35 + j * 0.55
        add_text(slide, num, 0.9, y, 0.4, 0.4, size=16, color=col, bold=True)
        add_text(slide, step, 1.4, y, 1.8, 0.4, size=16, color=WHITE, bold=True)
        add_text(slide, desc, 3.3, y, 2.8, 0.4, size=14, color=GRAY)

    # Outils
    add_card(slide, 7.0, 1.5, 5.5, 5.0, CARD_BG, GREEN, 1.5)
    add_text(slide, "OUTILS & INFRASTRUCTURE", 7.3, 1.7, 4.9, 0.35,
             size=17, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_line_h(slide, 7.3, 2.15, 4.9, BORDER, 1)

    tools = [
        ("Git + GitHub Actions", "Versioning, CI/CD, quality gates"),
        ("Docker Compose", "4 conteneurs, déploiement 1 commande"),
        ("MongoDB", "Stockage 245K posts, indexes, validation"),
        ("FastAPI", "API REST pour intégration externe"),
        ("pytest + mutmut", "501 tests, 80% coverage, 80.3% kill"),
        ("CodeCarbon", "Suivi empreinte carbone par version"),
        ("Gantt 16 WP", "28 jalons, planification détaillée"),
    ]
    for j, (tool, desc) in enumerate(tools):
        y = 2.35 + j * 0.55
        add_text(slide, tool, 7.4, y, 2.5, 0.4, size=15, color=WHITE, bold=True)
        add_text(slide, desc, 10.0, y, 2.3, 0.4, size=13, color=GRAY)

    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 13 \u2014 MÉTHODOLOGIE (15:00 \u2192 16:00)
Timing : 60 secondes (Sébastien)

Texte voix-off (Sébastien) :
\u00ab Un mot sur notre méthodologie.
On a suivi le cycle CRISP-DM adapté au machine learning :
comprendre, explorer, préparer, modéliser, évaluer, déployer.
Et surtout : itérer. 9 versions en 6 mois.

La gestion de projet s'est appuyée sur un planning Gantt
avec 16 work packages et 28 jalons.
Nos outils : Git avec CI/CD sur GitHub Actions,
Docker Compose pour le déploiement, MongoDB pour le stockage,
FastAPI pour l'API REST, et CodeCarbon pour le suivi carbone.
Le tout versionné et reproductible. \u00bb""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 14 — ROI ET VALEUR BUSINESS                         15:45
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "ROI & BUDGET — NIAMATO CONSULTING", YELLOW)

    # ── Panneau gauche : ROI (4 mini-cartes) ──
    add_card(slide, 0.4, 1.4, 6.4, 5.4, CARD_BG, YELLOW, 1.5)
    add_text(slide, "RETOUR SUR INVESTISSEMENT", 0.6, 1.55, 6.0, 0.35,
             size=16, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)
    add_line_h(slide, 0.8, 2.0, 5.6, BORDER, 1)

    roi_mini = [
        ("TEMPS",   "x10",  "5 min → 30s / post",      CYAN),
        ("VOLUME",  "x200", "60K posts/jour auto.",      GREEN),
        ("RISQUE",  "-67%", "faux positifs éliminés",    YELLOW),
        ("COÛT",    "≈ 0",  "0,0005 ct / post analysé",  RED),
    ]
    for j, (label, big, desc, col) in enumerate(roi_mini):
        y = 2.15 + j * 1.15
        add_text(slide, label, 0.7, y, 1.2, 0.3, size=14, color=col, bold=True)
        add_text(slide, big, 2.0, y - 0.1, 1.8, 0.5, size=36, color=col, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, desc, 3.9, y + 0.02, 2.7, 0.3, size=14, color=GRAY)

    # ── Panneau droit : Budget ──
    add_card(slide, 7.2, 1.4, 5.3, 5.4, CARD_BG, CYAN, 1.5)
    add_text(slide, "BUDGET PROJET", 7.4, 1.55, 4.9, 0.35,
             size=16, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_line_h(slide, 7.5, 2.0, 4.7, BORDER, 1)

    budget_lines = [
        ("Ressources humaines",  "49 500 €", WHITE),
        ("  Lead Technique (65 j)", "29 250 €", GRAY),
        ("  Validation ML (45 j)",  "20 250 €", GRAY),
        ("Infrastructure",       "0 €",      WHITE),
        ("Licences / outils",    "0 €",      WHITE),
        ("Annotation gold set",  "750 €",    WHITE),
    ]
    y = 2.2
    for label, val, col in budget_lines:
        indent = 0.2 if label.startswith("  ") else 0.0
        sz = 13 if label.startswith("  ") else 15
        add_text(slide, label, 7.6 + indent, y, 3.0, 0.3, size=sz, color=col)
        add_text(slide, val, 10.5, y, 1.8, 0.3, size=sz, color=col, bold=True,
                 align=PP_ALIGN.RIGHT)
        y += 0.35

    add_line_h(slide, 7.6, y + 0.05, 4.5, CYAN, 2)
    y += 0.2
    add_text(slide, "TOTAL PROJET", 7.6, y, 3.0, 0.35, size=17, color=CYAN, bold=True)
    add_text(slide, "50 250 €", 10.5, y, 1.8, 0.35, size=17, color=CYAN, bold=True,
             align=PP_ALIGN.RIGHT)

    y += 0.55
    add_line_h(slide, 7.6, y, 4.5, BORDER, 1)
    y += 0.15
    add_text(slide, "Exploitation mensuelle", 7.6, y, 3.0, 0.3, size=14, color=GREEN)
    add_text(slide, "~930 €/mois", 10.5, y, 1.8, 0.3, size=14, color=GREEN, bold=True,
             align=PP_ALIGN.RIGHT)
    y += 0.35
    add_text(slide, "100% open source\nAucun coût de licence", 7.6, y, 4.5, 0.5,
             size=12, color=DARK_GRAY, italic=True)

    add_text(slide, "TJM consultant junior : 450 €  |  Stack : Python, PyTorch, Streamlit, Docker (OSS)",
             0.4, 6.85, 12.3, 0.3, size=12, color=DARK_GRAY, italic=True,
             align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 14 — ROI & BUDGET (15:45 → 16:30)
Speaker : SÉBASTIEN
Timing : 45 secondes

Texte voix-off (Sébastien) :
« Parlons chiffres. Le projet Thumalien a coûté environ 50 000 euros,
essentiellement en ressources humaines — 110 jours-homme répartis
entre le développement technique et la validation.

Zéro euro de licence : notre stack est 100 % open source.
Zéro euro de cloud : tout tourne en local sur Docker.
Le seul coût additionnel : 750 euros pour l'annotation manuelle
du gold set de validation.

En exploitation, le coût mensuel est d'environ 930 euros :
un petit serveur à 30 euros et 2 jours de maintenance.
Pour 1,8 million de posts analysés par mois,
ça revient à 0,0005 centime par post.

En retour : un gain de productivité x10 pour les modérateurs,
une couverture de 60 000 posts par jour — 200 fois plus qu'un humain —
et une réduction de 67 % des faux positifs. »""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 15 — ROADMAP V10-V12                                16:30
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "ROADMAP — VERS V12")
    add_subtitle(slide, "Limites assumées et perspectives")

    roadmap = [
        ("V10", "Q3 2026", [
            "MLflow tracking",
            "Drift monitoring",
            "Streamlit AppTest",
        ], "Industrialisation", CYAN),
        ("V11", "Q4 2026", [
            "ClaimBuster intégration",
            "Monitoring Grafana",
            "Tests E2E nightly",
        ], "Vérif. factuelle", GREEN),
        ("V12", "2027", [
            "Modèle FR spécialisé",
            "Annotation communautaire",
            "Federated learning",
        ], "Inclusion", YELLOW),
    ]
    for i, (ver, date, items, theme, col) in enumerate(roadmap):
        x = 0.5 + i * 4.2
        add_card(slide, x, 1.5, 3.8, 5.2, CARD_BG, col, 1.5)
        add_text(slide, ver, x + 0.1, 1.7, 3.6, 0.6,
                 size=40, color=col, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, date, x + 0.1, 2.35, 3.6, 0.35,
                 size=20, color=GRAY, align=PP_ALIGN.CENTER)
        add_line_h(slide, x + 0.3, 2.85, 3.2, BORDER, 1)

        for j, item in enumerate(items):
            add_text(slide, f"•  {item}", x + 0.4, 3.1 + j * 0.6, 3.2, 0.5,
                     size=17, color=WHITE)

        add_line_h(slide, x + 0.3, 5.6, 3.2, BORDER, 1)
        add_text(slide, theme, x + 0.1, 5.75, 3.6, 0.4,
                 size=16, color=YELLOW, bold=True, align=PP_ALIGN.CENTER,
                 italic=True)

    add_footer(slide)
    add_slide_number(slide, n)

    set_notes(slide, """SLIDE 12 — ROADMAP (14:30 → 16:30)
Timing : 2 minutes

Texte voix-off :
« Nous sommes lucides sur les limites de V9.
La roadmap est structurée en 3 phases.
V10 au troisième trimestre 2026 : industrialisation avec MLflow pour le tracking
d'expériences et monitoring de drift.
V11 fin 2026 : intégration de ClaimBuster pour la vérification factuelle,
monitoring Grafana, et tests end-to-end automatisés chaque nuit.
V12 en 2027 : un modèle spécialisé français, de l'annotation communautaire,
et du federated learning pour l'inclusion.
Chaque version a un thème clair et des livrables concrets. »

Point clé : montrer qu'on sait où on va, sans sur-promettre.""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 13a — CITATION FINALE                               17:30
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BLACK)

    # Citation centrée
    add_text(slide,
             "« Un score sans explication\n  est un verdict sans procès. »",
             1.0, 2.2, 11.3, 3.0,
             size=56, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    # Filet décoratif
    add_line_h(slide, 4.5, 5.8, 4.3, CYAN_DIM, 1)

    set_notes(slide, """SLIDE 13 — CITATION (17:30 → 17:45)
Timing : 15 secondes — SILENCE après avoir prononcé la phrase

Animation : la citation reste 5 secondes complètes en silence,
puis fondu enchaîné de 1.5s vers la slide remerciements.

Texte voix-off (prononcé lentement, regard caméra) :
« Un score sans explication est un verdict sans procès. »
[SILENCE 5 secondes — laisser résonner]

Ce silence renforce l'impact. Ne rien ajouter.""")

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 13b — REMERCIEMENTS                                 17:45
    # ═════════════════════════════════════════════════════════════════
    n += 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)

    add_text(slide, "MERCI", 0.5, 1.3, 12.3, 1.1,
             size=72, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    add_line_h(slide, 4.0, 2.6, 5.3, CYAN, 2)

    add_text(slide, "Azélie Bernard  ·  Sébastien Lazcanotegui",
             0.5, 2.9, 12.3, 0.5, size=24, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Master 1 Big Data & Intelligence Artificielle — 2026",
             0.5, 3.5, 12.3, 0.4, size=18, color=GRAY, align=PP_ALIGN.CENTER)

    add_line_h(slide, 1.5, 4.3, 10.3, BORDER, 1)

    # QR placeholders
    qr_items = [
        ("Repository GitHub", 2.0),
        ("Rapport complet",   5.5),
        ("Model Card",        9.0),
    ]
    for label, x in qr_items:
        add_card(slide, x, 4.6, 2.3, 1.6, CARD_BG, CYAN_DIM, 1)
        add_text(slide, "QR", x, 4.7, 2.3, 1.0,
                 size=28, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, 6.3, 2.3, 0.35,
                 size=14, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, "Questions ?  Disponible dans le chat de soutenance.",
             0.5, 7.0, 12.3, 0.35, size=17, color=GRAY, italic=True,
             align=PP_ALIGN.CENTER)

    add_footer(slide)

    set_notes(slide, """SLIDE 14 — REMERCIEMENTS (17:45 → 18:00)
Timing : 15 secondes

Texte voix-off :
« Merci de votre attention. Le repository, le rapport complet
et la Model Card sont accessibles via ces QR codes.
Nous sommes disponibles pour vos questions. »

Rappel : insérer les vrais QR codes avant export final :
- QR 1 → URL repo GitHub
- QR 2 → docs/pdf/rapport_projet_thumalien.pdf
- QR 3 → docs/12_model_card.md""")

    return prs


# ╔═══════════════════════════════════════════════════════════════════╗
# ║                          MAIN                                    ║
# ╚═══════════════════════════════════════════════════════════════════╝

if __name__ == "__main__":
    prs = create_presentation()
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "Thumalien_V9_slides.pptx")
    prs.save(out_path)
    print(f"✅ Présentation sauvegardée : {out_path}")
    print(f"   {len(prs.slides)} slides générées")
    print()
    print("📋 Checklist avant export final :")
    print("   1. Insérer le logo Thumalien (slide 1)")
    print("   2. Insérer les 3 QR codes (slide 14)")
    print("   3. Vérifier la police Inter (installer si absente)")
    print("   4. Tester la lisibilité sur téléphone à 50 cm")
    print("   5. Exporter en PNG 1920×1080 pour montage DaVinci")
