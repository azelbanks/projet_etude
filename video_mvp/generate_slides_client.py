#!/usr/bin/env python3
"""
Deck CLIENT — ThumaCheck (Niamato Consulting pour Thumalien)
Design aligne sur la charte des logos Niamato : bleu petrole #0E3F40,
corail #EF8265, beige #F2E7D6. Esthetique consulting claire et professionnelle.
Genere ThumaCheck_slides_CLIENT.pptx — n'ecrase NI le pptx NI le script existants.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy

# ---- Palette marque Niamato ----
PETROL   = RGBColor(0x0E,0x3F,0x40)
PETROL_D = RGBColor(0x09,0x2B,0x2C)   # fond sombre plus profond
CORAL    = RGBColor(0xEF,0x82,0x65)
CORAL_D  = RGBColor(0xD9,0x6A,0x4D)
BEIGE    = RGBColor(0xF2,0xE7,0xD6)
PAGE_BG  = RGBColor(0xFB,0xF9,0xF4)   # blanc chaud (fond slides claires)
CARD     = RGBColor(0xFF,0xFF,0xFF)
CARD_ALT = RGBColor(0xF6,0xF1,0xE8)
BORDER   = RGBColor(0xE4,0xDC,0xCD)
BLUEGREY = RGBColor(0xB0,0xC1,0xC0)
LIGHTBL  = RGBColor(0x95,0xB5,0xE8)
GREEN    = RGBColor(0x24,0x75,0x49)
INK      = RGBColor(0x20,0x2E,0x2F)   # texte principal
MUTE     = RGBColor(0x7C,0x85,0x84)   # texte secondaire
WHITE    = RGBColor(0xFF,0xFF,0xFF)
BEIGE_TX = RGBColor(0xD9,0xCF,0xBE)   # texte clair sur fond sombre

FONT = "Inter"
EMUIN = 914400
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
OUT = "/Users/azeliebernard/Documents/MASTER Big data/projet_etude/video_mvp/ThumaCheck_slides_CLIENT.pptx"
ASSET = "/Users/azeliebernard/Desktop/thumalien_captures/qr_codes/"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# =================== HELPERS ===================
def slide(bg=PAGE_BG):
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = bg
    return s

def no_line(sh):
    sh.line.fill.background()

def no_shadow(sh):
    el = sh._element.spPr
    # explicit no shadow
    return

def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def line(s, x, y, w, h, color=CORAL, lw=2.0):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=0, line_spacing=None, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after: p.space_after = Pt(space_after)
        if line_spacing: p.line_spacing = line_spacing
        for (t, sz, col, bold, *rest) in para:
            it = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.name = FONT
            r.font.color.rgb = col; r.font.bold = bold; r.font.italic = it
    return tb

def P(*runs):  # convenience single paragraph
    return list(runs)

def pic(s, path, x, y, w=None, h=None):
    kw = {}
    if w is not None: kw["width"] = Inches(w)
    if h is not None: kw["height"] = Inches(h)
    return s.shapes.add_picture(path, Inches(x), Inches(y), **kw)

def footer(s, n):
    line(s, 0.7, 7.0, 11.93, 0.011, color=BORDER)
    text(s, 0.7, 7.06, 9.0, 0.3,
         [P(("Niamato Consulting   ·   ThumaCheck   ·   Client : Thumalien   ·   M1 BDIA 2026", 9, MUTE, False))],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 11.0, 7.06, 1.63, 0.3, [P((f"{n} / 17", 9, MUTE, True))],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # petit logo bleu en haut a droite
    try:
        pic(s, ASSET+"logo_h_blue.png", 10.55, 0.5, w=2.05)
    except Exception: pass

def head(s, kicker, title, title_color=PETROL):
    text(s, 0.7, 0.55, 8.0, 0.3, [P((kicker.upper(), 12, CORAL, True))])
    text(s, 0.7, 0.86, 9.6, 0.95, [P((title, 33, title_color, True))], line_spacing=1.0)
    line(s, 0.72, 1.62, 0.95, 0.05, color=CORAL)

# =================== SLIDE 1 — COUVERTURE ===================
s = slide(PETROL)
# bande coral en bas
line(s, 0, 7.30, 13.333, 0.20, color=CORAL)
try: pic(s, ASSET+"logo_h_white.png", 4.5, 1.05, w=4.33)
except Exception: pass
text(s, 1.0, 2.75, 11.33, 1.1, [P(("THUMACHECK", 60, WHITE, True))], align=PP_ALIGN.CENTER)
text(s, 1.0, 3.95, 11.33, 0.5,
     [P(("Détection de désinformation sur Bluesky", 21, BEIGE_TX, False))], align=PP_ALIGN.CENTER)
text(s, 1.0, 4.40, 11.33, 0.4,
     [P(("bilingue   ·   explicable   ·   frugale", 15, CORAL, True))], align=PP_ALIGN.CENTER)
line(s, 5.67, 5.10, 2.0, 0.018, color=BLUEGREY)
text(s, 1.0, 5.35, 11.33, 0.4,
     [P(("Azélie Bernard   ·   Sébastien Lazcanotegui", 17, WHITE, True))], align=PP_ALIGN.CENTER)
text(s, 1.0, 5.78, 11.33, 0.35,
     [P(("Master 1 Big Data & Intelligence Artificielle — 2025 / 2026", 13, BEIGE_TX, False))], align=PP_ALIGN.CENTER)
text(s, 1.0, 6.45, 11.33, 0.3,
     [P(("Niamato Consulting  pour  Thumalien", 12, BLUEGREY, False))], align=PP_ALIGN.CENTER)

# =================== SLIDE 2 — EQUIPE ===================
s = slide()
head(s, "Niamato Consulting", "L'équipe projet")
def member(x, name, role, bullets, accent, initials):
    rect(s, x, 1.95, 5.55, 4.55, fill=CARD, line=BORDER, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    rect(s, x, 1.95, 5.55, 0.10, fill=accent)
    # pastille photo (placeholder initiales)
    rect(s, x+0.35, 2.35, 1.0, 1.0, fill=CARD_ALT, line=accent, lw=1.5, shape=MSO_SHAPE.OVAL)
    text(s, x+0.35, 2.35, 1.0, 1.0, [P((initials, 26, accent if accent!=BEIGE else CORAL_D, True))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+1.55, 2.42, 3.7, 0.5, [P((name, 19, PETROL, True))])
    text(s, x+1.55, 2.86, 3.7, 0.4, [P((role, 13, accent if accent!=BEIGE else CORAL_D, True))])
    runs = [P(("•  "+b, 12.5, INK, False)) for b in bullets]
    text(s, x+0.40, 3.65, 4.85, 2.7, runs, space_after=6, line_spacing=1.0)
member(0.7, "Azélie Bernard", "Lead Technique",
        ["Pipeline NLP V1 → V9", "Dashboard Streamlit 5 pages",
         "XAI : SHAP, Captum, faithfulness", "Fine-tuning CamemBERT & RoBERTa",
         "CI/CD, tests, mutation testing"], PETROL, "AB")
member(7.08, "Sébastien Lazcanotegui", "Validation & Qualité ML",
        ["Annotation gold set (473 posts)", "Débiaisage Reuters (BODY_AGENCY_TERMS)",
         "GridSearch hyperparamètres", "Revue critique & challenge technique",
         "Conformité RGPD & AI Act"], CORAL, "SL")
footer(s, 2)

# =================== SLIDE 3 — PROBLEME CHIFFRE ===================
s = slide()
head(s, "Le constat", "La désinformation sur Bluesky")
stats = [("35 M","utilisateurs Bluesky",PETROL),
         ("60 000+","posts publics FR / EN par jour",CORAL),
         ("0","équipe de modération centralisée",PETROL),
         ("Angle mort","de la régulation européenne",CORAL)]
gx=[0.7,6.94]; gy=[2.05,4.30]; cw=5.69; ch=2.05
for i,(big,leg,acc) in enumerate(stats):
    x=gx[i%2]; y=gy[i//2]
    rect(s,x,y,cw,ch,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
    rect(s,x,y,0.10,ch,fill=acc)
    sz = 58 if big not in ("Angle mort",) else 40
    text(s,x+0.5,y+0.32,cw-0.8,1.15,[P((big,sz,acc,True))],anchor=MSO_ANCHOR.MIDDLE)
    text(s,x+0.5,y+1.45,cw-0.8,0.5,[P((leg,15,INK,False))])
footer(s,3)

# =================== SLIDE 4 — CAHIER DES CHARGES ===================
s = slide()
head(s, "Le mandat", "Le cahier des charges")
reqs=[("TRANSPARENCE","Chaque score doit être explicable et auditable",PETROL),
      ("BILINGUE","Français + anglais, performance équivalente",LIGHTBL),
      ("FRUGALITÉ","< 5 ms par texte, empreinte CO₂ mesurée",GREEN),
      ("CONFORMITÉ","RGPD art. 22 + AI Act art. 13 / 14 / 50 — 2 août 2026",CORAL)]
y=1.95
for i,(k,d,acc) in enumerate(reqs):
    yy=y+i*1.12
    rect(s,0.7,yy,11.93,0.98,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.10)
    rect(s,0.7,yy,0.14,0.98,fill=acc)
    rect(s,1.05,yy+0.19,0.6,0.6,fill=acc,shape=MSO_SHAPE.OVAL)
    text(s,1.05,yy+0.19,0.6,0.6,[P((str(i+1),20,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,2.0,yy,3.4,0.98,[P((k,21,PETROL,True))],anchor=MSO_ANCHOR.MIDDLE)
    text(s,5.3,yy,7.1,0.98,[P((d,16,INK,False))],anchor=MSO_ANCHOR.MIDDLE)
text(s,0.7,6.55,11.93,0.35,
     [P(("Quatre exigences non négociables pour passer du prototype au système déployable.",13,MUTE,False,True))])
footer(s,4)

# =================== SLIDE 5 — LE PIEGE F1 ===================
s = slide()
head(s, "Le tournant", "F1 = 0.99 — le piège", title_color=CORAL_D)
# gauche : metriques vertes
rect(s,0.7,1.95,5.7,4.05,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
text(s,1.05,2.2,5.0,0.4,[P(("CROSS-VALIDATION 5-FOLD",13,PETROL,True))])
mets=[("F1 macro","0.987"),("Précision","0.992"),("Recall","0.984"),("Accuracy","0.989")]
for i,(k,v) in enumerate(mets):
    yy=2.75+i*0.62
    text(s,1.05,yy,3.0,0.4,[P((k,15,INK,False))])
    text(s,4.3,yy,1.8,0.4,[P((v,17,GREEN,True))],align=PP_ALIGN.RIGHT)
line(s,1.05,5.35,4.95,0.014,color=BORDER)
text(s,1.05,5.5,5.0,0.4,[P(("✓  Tous les voyants au vert",13,GREEN,True))])
# droite : XAI revelation
rect(s,6.93,1.95,5.7,4.05,fill=PETROL,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
text(s,7.28,2.2,5.0,0.4,[P(("TOP COEFFICIENTS LOGREG (XAI)",13,CORAL,True))])
coef=[("reuters","+0.84"),("afp","+0.71"),("ap_news","+0.69"),("associated","+0.63"),("press","+0.58")]
for i,(k,v) in enumerate(coef):
    yy=2.72+i*0.45
    text(s,7.28,yy,3.0,0.35,[P((k,14,WHITE,False))])
    text(s,10.5,yy,1.8,0.35,[P((v,14,CORAL,True))],align=PP_ALIGN.RIGHT)
text(s,7.28,5.15,5.05,0.8,[P(("▲  Le modèle apprend le STYLE des",13,CORAL,True)),
                            P(("     agences de presse, pas les faits.",13,BEIGE_TX,True))],line_spacing=1.05)
text(s,0.7,6.45,11.93,0.4,
     [P(("Sans XAI, nous aurions livré un modèle qui s'effondre en production.",14,CORAL_D,True,True))])
footer(s,5)

# =================== SLIDE 6 — ARCHITECTURE ===================
s = slide()
head(s, "La solution", "Architecture — 8 conteneurs")
boxes=[("Bluesky","AT Proto",BLUEGREY),("Collector","Python",PETROL),("MongoDB","storage",PETROL),
       ("Detector","V5 + V6 + CamBERT",CORAL),("Pipeline","Cascade Stage 1+2",PETROL),("Streamlit","Dashboard",PETROL),
       ("XAI Engine","SHAP + Captum",CORAL),("Monitoring","CodeCarbon",GREEN)]
bw=2.75; bh=1.15; gxs=[0.7,3.75,6.8,9.85]
def cbox(x,y,t,sub,acc):
    rect(s,x,y,bw,bh,fill=CARD,line=acc,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    rect(s,x,y,bw,0.09,fill=acc)
    text(s,x+0.15,y+0.22,bw-0.3,0.45,[P((t,16,PETROL,True))],align=PP_ALIGN.CENTER)
    text(s,x+0.15,y+0.66,bw-0.3,0.4,[P((sub,11.5,MUTE,False))],align=PP_ALIGN.CENTER)
row1_y=2.15; row2_y=4.35
for i in range(4): cbox(gxs[i],row1_y,*boxes[i])
for i in range(4): cbox(gxs[i],row2_y,*boxes[i+4])
# fleches horizontales rangee 1
for i in range(3):
    ax=gxs[i]+bw+0.04
    text(s,ax,row1_y+0.2,0.25,bh,[P(("›",24,CORAL,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# fleches verticales
for i in range(4):
    text(s,gxs[i],row1_y+bh+0.02,bw,0.5,[P(("⌄",20,CORAL,True))],align=PP_ALIGN.CENTER)
text(s,0.7,6.55,11.93,0.35,
     [P(("Données → IA → Décision   |   C4 Model   |   Docker Compose   |   FastAPI",13,MUTE,False)) ],align=PP_ALIGN.CENTER)
footer(s,6)

# =================== SLIDE 7 — PIPELINE CASCADE ===================
s = slide()
head(s, "Sous le capot", "Pipeline cascade — inférence")
actors=["User","Stage 1","V5","V6","CamBERT","V8 Méta","XAI"]
n=len(actors); x0=1.0; xw=(11.93-0.0)/n
for i,a in enumerate(actors):
    cx=x0+i*( (12.0-x0) /(n-0.0))
for i,a in enumerate(actors):
    cx=0.95+i*1.70
    acc = CORAL if a in ("Stage 1","XAI") else PETROL
    rect(s,cx,2.05,1.45,0.62,fill=acc,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.18)
    text(s,cx,2.05,1.45,0.62,[P((a,12.5,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    # ligne de vie
    line(s,cx+0.72,2.70,0.012,3.1,color=BORDER)
msgs=[("texte",0,1),("opinion ?",1,1),("si factuel",1,2),("score TF-IDF",2,3),
      ("score style",3,4),("embedding",4,5),("décomposition",5,6),("score + XAI",6,0)]
for j,(m,a,b) in enumerate(msgs):
    yy=3.0+j*0.42
    xa=0.95+a*1.70+0.72; xb=0.95+b*1.70+0.72
    if a==b:
        line(s,xa,yy+0.16,0.85,0.016,color=CORAL)
        text(s,xa+0.1,yy-0.08,2.2,0.3,[P((m,11,PETROL,True))],align=PP_ALIGN.LEFT,wrap=False)
    else:
        lo,hi=sorted([xa,xb])
        line(s,lo,yy+0.16,hi-lo,0.016,color=CORAL)
        text(s,lo,yy-0.06,hi-lo,0.3,[P((m,11,PETROL,True))],align=PP_ALIGN.CENTER)
rect(s,0.7,6.45,3.2,0.5,fill=CARD_ALT,line=CORAL,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.2)
text(s,0.7,6.45,3.2,0.5,[P(("▸  1.5 ms au total",14,PETROL,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
rect(s,9.4,6.45,3.23,0.5,fill=CARD_ALT,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.2)
text(s,9.4,6.45,3.23,0.5,[P(("opinion  ≠  fake news",13,PETROL,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
footer(s,7)

# =================== SLIDE 8 — FAITHFULNESS METHODE ===================
s = slide()
head(s, "La rigueur", "La fidélité des explications")
text(s,0.7,1.74,9.6,0.4,[P(("Au-delà de SHAP : la mesurer",15,MUTE,False,True))])
# question
rect(s,0.7,2.15,11.93,1.05,fill=PETROL,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
text(s,1.1,2.15,11.1,1.05,
     [P(("Comment savoir que vos explications reflètent vraiment",16,WHITE,True)),
      P(("le comportement réel du modèle ?",16,WHITE,True))],
     anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.1)
text(s,0.7,3.30,11.93,0.4,[P(("⌄",24,CORAL,True))],align=PP_ALIGN.CENTER)
# methode
rect(s,0.7,3.78,11.93,2.55,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
text(s,1.1,4.0,11.1,0.4,[P(("PROTOCOLE ERASER  ",14,CORAL,True),("(DeYoung et al., ACL 2020)",13,MUTE,False))])
steps=["Identifier les top-k features via SHAP","Masquer ces features (valeur = 0)",
       "Mesurer la chute de P(suspect)","Comparer à un masquage aléatoire"]
for i,st in enumerate(steps):
    yy=4.55+i*0.40
    text(s,1.3,yy,0.5,0.35,[P((f"{i+1}.",14,CORAL,True))])
    text(s,1.85,yy,10.0,0.35,[P((st,14,INK,False))])
text(s,1.1,6.22,11.1,0.4,
     [P(("Si l'explication est fidèle :   ",13,PETROL,True),("chute_attribution  ≫  chute_random",13,GREEN,True))])
footer(s,8)

# =================== SLIDE 9 — RESULTAT AOPC ===================
s = slide()
head(s, "Le résultat", "Faithfulness — uplift + 21 %", title_color=GREEN)
# grand chiffre gauche
rect(s,0.7,1.95,4.3,4.05,fill=PETROL,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
text(s,0.9,2.7,3.9,1.2,[P(("+0.21",62,CORAL,True))],align=PP_ALIGN.CENTER)
text(s,0.9,3.95,3.9,0.4,[P(("uplift d'attribution",15,WHITE,False))],align=PP_ALIGN.CENTER)
text(s,0.9,4.55,3.9,0.5,[P(("5.6×",30,WHITE,True))],align=PP_ALIGN.CENTER)
text(s,0.9,5.15,3.9,0.4,[P(("vs masquage aléatoire",13,BEIGE_TX,False))],align=PP_ALIGN.CENTER)
# tableau droite
rect(s,5.3,1.95,7.33,4.05,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
text(s,5.7,2.2,5.0,0.4,[P(("MÉTRIQUES",13,PETROL,True))])
rows=[("AOPC attribution","0.253",GREEN),("AOPC random","0.045",MUTE),
      ("Comprehensiveness @5","0.232",PETROL),("Sufficiency @5","0.058",PETROL)]
for i,(k,v,c) in enumerate(rows):
    yy=2.75+i*0.62
    text(s,5.7,yy,5.0,0.4,[P((k,15,INK,False))])
    text(s,10.6,yy,1.65,0.4,[P((v,18,c,True))],align=PP_ALIGN.RIGHT)
line(s,5.7,5.3,6.55,0.014,color=BORDER)
text(s,5.7,5.45,6.6,0.4,[P(("✓  Cible > +0.10   —   atteinte ×2",14,GREEN,True))])
footer(s,9)

# =================== SLIDE 10 — QUALITE INDUSTRIELLE ===================
s = slide()
head(s, "L'exigence", "Qualité industrielle")
text(s,0.7,1.78,11.0,0.4,[P(("537 tests   ·   80 % coverage   ·   80,3 % mutation kill rate",15,MUTE,False,True))])
panels=[("PYTEST","537","tests passants",["0 échec","0 skip"],PETROL),
        ("COVERAGE","80 %","line coverage",["77,9 % branch","quality gate CI --fail-under=75"],CORAL),
        ("MUTMUT","80,3 %","kill rate (143 / 178)",["> benchmark Google","(60–75 %)"],GREEN)]
pw=3.85; px=[0.7,4.74,8.78]
for i,(t,big,sub,extra,acc) in enumerate(panels):
    x=px[i]
    rect(s,x,2.15,pw,4.1,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
    rect(s,x,2.15,pw,0.10,fill=acc)
    text(s,x,2.45,pw,0.4,[P((t,14,MUTE,True))],align=PP_ALIGN.CENTER)
    text(s,x,2.95,pw,1.1,[P((big,52,acc,True))],align=PP_ALIGN.CENTER)
    text(s,x,4.05,pw,0.4,[P((sub,14,INK,True))],align=PP_ALIGN.CENTER)
    runs=[P((e,12,MUTE,False)) for e in extra]
    text(s,x+0.3,4.75,pw-0.6,1.2,runs,align=PP_ALIGN.CENTER,space_after=4)
text(s,0.7,6.5,11.93,0.35,[P(("De la quantité à la qualité — chaque ligne de production est validée.",13,MUTE,False,True))])
footer(s,10)

# =================== SLIDE 11 — CONFORMITE ===================
s = slide()
head(s, "Le cadre", "Conformité réglementaire")
def reg_block(x, title, sub, items, acc):
    rect(s,x,2.0,5.83,4.8,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
    rect(s,x,2.0,5.83,0.10,fill=acc)
    text(s,x+0.35,2.28,5.2,0.4,[P((title,16,PETROL,True))])
    text(s,x+0.35,2.68,5.2,0.35,[P((sub,12,MUTE,False,True))])
    for i,(art,prv) in enumerate(items):
        yy=3.2+i*0.80
        text(s,x+0.35,yy,0.45,0.4,[P(("✓",16,GREEN,True))])
        text(s,x+0.85,yy,4.7,0.4,[P((art,14,INK,True))])
        text(s,x+0.85,yy+0.36,4.7,0.4,[P((prv,12,MUTE,False))])
reg_block(0.7,"AI Act  (UE 2024/1689)","pleinement applicable 2 août 2026",
          [("Art. 13 — transparence","Model Card MC-THUM-2026-001"),
           ("Art. 14 — supervision","Décomposition β·x dans le dashboard"),
           ("Art. 50 — transparence IA","Bannière IA visible dans le dashboard"),
           ("Risque limité classifié","Doc 02 § 4.1")],CORAL)
reg_block(6.80,"RGPD  (UE 2016/679)","protection des données personnelles",
          [("Art. 22 — décision auto","Droit à l'explication SHAP + Captum"),
           ("Art. 35 — AIPD","Document RGPD-THUM-2026-001"),
           ("Base légale art. 6.1.f","Posts publics, intérêt légitime")],PETROL)
text(s,0.7,6.95,11.93,0.35,[P(("Positionnement FLI AI Safety Index : supervision humaine, explications par défaut, empreinte mesurée",12,MUTE,False,True))])
footer(s,11)

# =================== SLIDE 12 — GREEN IT ===================
s = slide()
head(s, "L'empreinte", "Green IT — 8,86 g CO₂ mesuré (CodeCarbon)", title_color=GREEN)
# repartition gauche
rect(s,0.7,1.95,6.0,4.05,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
text(s,1.05,2.2,5.3,0.4,[P(("RÉPARTITION DE L'EMPREINTE",13,PETROL,True))])
parts=[("RoBERTa EN (V1+V2)","39 %  (3,47 g)",0.39,PETROL),
       ("Expert V3+V4+V5","43 %  (3,78 g)",0.43,CORAL),
       ("V6 Style GradientBoost","8 %  (0,69 g)",0.08,BLUEGREY),
       ("V5 Bluesky fine-tune","5 %  (0,45 g)",0.05,GREEN),
       ("CamemBERT V1+V2","5 %  (0,45 g)",0.05,LIGHTBL),
       ("V7 + V8 + Stage1","< 1 %  (0,01 g)",0.01,MUTE)]
for i,(k,v,frac,c) in enumerate(parts):
    yy=2.72+i*0.52
    text(s,1.05,yy,3.1,0.35,[P((k,12.5,INK,False))])
    text(s,5.0,yy,1.55,0.35,[P((v,12,MUTE,True))],align=PP_ALIGN.RIGHT)
    rect(s,1.05,yy+0.30,5.45*max(frac,0.012),0.09,fill=c)
# decision droite
rect(s,6.93,1.95,5.7,4.05,fill=CARD_ALT,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
text(s,7.28,2.2,5.0,0.4,[P(("DÉCISION ARCHITECTURALE",13,PETROL,True))])
dec=[("Production temps réel","V5 seul — 1,5 ms, 0,6 g CO₂ / jour"),
     ("Analyse offline batch","V8 avec CamemBERT — 0,45 g entraînement"),
     ("ROI Green IT","frugalité en prod + puissance en analyse")]
for i,(k,v) in enumerate(dec):
    yy=2.75+i*1.0
    text(s,7.28,yy,5.0,0.35,[P(("→  "+k,14,CORAL_D,True))])
    text(s,7.55,yy+0.36,4.8,0.5,[P((v,12.5,INK,False))])
text(s,7.28,5.75,5.1,0.4,[P(("Documenté — Model Card § 8",11.5,MUTE,False,True))])
text(s,0.7,6.5,11.93,0.35,
     [P(("Mesuré via CodeCarbon (Apple M4 Pro, France).   ≈ 52 m en voiture essence — réf. ADEME 2024 : 170 g CO₂/km — empreinte quasi nulle.",12,MUTE,False,True))])
footer(s,12)

# =================== SLIDE 13 — METHODOLOGIE ===================
s = slide()
head(s, "La démarche", "Méthodologie & organisation")
text(s,0.7,1.78,11.0,0.4,[P(("CRISP-DM adapté   ·   9 versions   ·   28 notebooks   ·   16 work packages",14,MUTE,False,True))])
steps=[("1","Comprendre","Cahier des charges, 4 exigences"),
       ("2","Explorer","245K posts, audit qualité"),
       ("3","Préparer","Débiaisage Reuters, synth."),
       ("4","Modéliser","V1→V9, LogReg + CamemBERT"),
       ("5","Évaluer","Gold set 473, faithfulness"),
       ("6","Déployer","Docker, CI/CD, dashboard"),
       ("9×","Itérer","9 versions en 6 mois")]
cw=1.62
for i,(num,t,d) in enumerate(steps):
    x=0.7+i*1.705
    acc = CORAL if num=="9×" else PETROL
    rect(s,x,2.15,cw,1.85,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    rect(s,x+0.51,2.32,0.6,0.6,fill=acc,shape=MSO_SHAPE.OVAL)
    text(s,x+0.51,2.32,0.6,0.6,[P((num,16,WHITE,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x+0.08,3.0,cw-0.16,0.35,[P((t,13,PETROL,True))],align=PP_ALIGN.CENTER)
    text(s,x+0.08,3.38,cw-0.16,0.55,[P((d,9.5,MUTE,False))],align=PP_ALIGN.CENTER)
    if i<len(steps)-1:
        text(s,x+cw-0.02,2.15,0.2,1.85,[P(("›",16,CORAL,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# outils
text(s,0.7,4.35,11.0,0.35,[P(("OUTILS & INFRASTRUCTURE",12,CORAL,True))])
tools=[("Git + GitHub Actions","Versioning, CI/CD, quality gates"),
       ("Docker Compose","4 conteneurs, déploiement 1 commande"),
       ("MongoDB","245K posts, indexes, validation"),
       ("FastAPI","API REST pour intégration externe"),
       ("pytest + mutmut","537 tests, 80 % cov, 80,3 % kill"),
       ("CodeCarbon","empreinte carbone par version")]
tw=3.85
for i,(t,d) in enumerate(tools):
    x=0.7+(i%3)*4.04; y=4.78+(i//3)*0.95
    rect(s,x,y,tw,0.82,fill=CARD_ALT,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    text(s,x+0.25,y+0.12,tw-0.4,0.35,[P((t,13,PETROL,True))])
    text(s,x+0.25,y+0.46,tw-0.4,0.3,[P((d,10.5,MUTE,False))])
footer(s,13)

# =================== SLIDE 14 — ROI & BUDGET ===================
s = slide()
head(s, "La valeur", "ROI & budget")
# KPI gauche
text(s,0.7,1.95,5.7,0.35,[P(("RETOUR SUR INVESTISSEMENT",13,CORAL,True))])
kpis=[("×10","5 min → 30 s par post"),("×200","60K posts/jour automatisés"),
      ("−67 %","faux positifs éliminés"),("≈ 0","0,0005 ct / post analysé")]
for i,(big,leg) in enumerate(kpis):
    x=0.7+(i%2)*2.95; y=2.4+(i//2)*1.65
    rect(s,x,y,2.75,1.45,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.06)
    text(s,x,y+0.18,2.75,0.7,[P((big,38,PETROL,True))],align=PP_ALIGN.CENTER)
    text(s,x+0.15,y+0.92,2.45,0.45,[P((leg,11.5,MUTE,False))],align=PP_ALIGN.CENTER)
# budget droite
rect(s,6.75,1.95,5.88,4.4,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.04)
text(s,7.1,2.2,5.2,0.4,[P(("BUDGET PROJET",13,PETROL,True))])
budget=[("Lead Technique (65 j)","29 250 €",INK),
        ("Validation ML (45 j)","20 250 €",INK),
        ("Infrastructure","0 €",GREEN),
        ("Licences / outils","0 €",GREEN),
        ("Annotation gold set","750 €",INK)]
for i,(k,v,c) in enumerate(budget):
    yy=2.7+i*0.50
    text(s,7.1,yy,3.8,0.35,[P((k,13,INK,False))])
    text(s,10.7,yy,1.6,0.35,[P((v,14,c,True))],align=PP_ALIGN.RIGHT)
line(s,7.1,5.28,5.18,0.016,color=BORDER)
text(s,7.1,5.42,3.8,0.4,[P(("TOTAL PROJET",14,PETROL,True))])
text(s,10.5,5.42,1.8,0.4,[P(("50 250 €",18,CORAL_D,True))],align=PP_ALIGN.RIGHT)
text(s,7.1,5.88,5.2,0.35,[P(("Exploitation : ~930 €/mois   ·   100 % open source",11.5,MUTE,False,True))])
text(s,0.7,6.55,11.93,0.35,
     [P(("TJM consultant junior : 450 €   |   Stack : Python, PyTorch, Streamlit, Docker (OSS)",12,MUTE,False,True))])
footer(s,14)

# =================== SLIDE 15 — ROADMAP ===================
s = slide()
head(s, "Et après", "Roadmap — vers V12")
text(s,0.7,1.78,11.0,0.4,[P(("Limites assumées et perspectives",15,MUTE,False,True))])
cols=[("V10","Q3 2026",["MLflow tracking","Drift monitoring","Streamlit AppTest"],"Industrialisation",PETROL),
      ("V11","Q4 2026",["ClaimBuster + LLM open-weight","Monitoring Grafana","Tests E2E nightly"],"Vérif. factuelle",CORAL),
      ("V12","2027",["Mistral / modèle souverain","Annotation communautaire","Federated learning"],"Inclusion",GREEN)]
cw=3.85
for i,(v,date,items,tag,acc) in enumerate(cols):
    x=0.7+i*4.04
    rect(s,x,2.2,cw,3.9,fill=CARD,line=BORDER,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
    rect(s,x,2.2,cw,0.75,fill=acc,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.05)
    rect(s,x,2.55,cw,0.40,fill=acc)
    text(s,x,2.28,cw,0.5,[P((v,26,WHITE,True))],align=PP_ALIGN.CENTER)
    text(s,x,3.05,cw,0.35,[P((date,14,MUTE,True))],align=PP_ALIGN.CENTER)
    for j,it in enumerate(items):
        text(s,x+0.4,3.55+j*0.55,cw-0.7,0.45,[P(("•  "+it,13.5,INK,False))])
    line(s,x+0.4,5.35,cw-0.8,0.014,color=BORDER)
    text(s,x,5.5,cw,0.45,[P((tag,14,acc,True))],align=PP_ALIGN.CENTER)
footer(s,15)

# =================== SLIDE 16 — CITATION ===================
s = slide(PETROL)
line(s,5.17,2.55,3.0,0.04,color=CORAL)
text(s,1.5,3.0,10.33,1.8,
     [P(("« Un score sans explication",40,WHITE,True,True)),
      P(("est un verdict sans procès. »",40,WHITE,True,True))],
     align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.15)
text(s,1.5,5.6,10.33,0.4,[P(("Niamato Consulting",13,BLUEGREY,True))],align=PP_ALIGN.CENTER)

# =================== SLIDE 17 — MERCI ===================
s = slide(PETROL)
line(s,0,7.30,13.333,0.20,color=CORAL)
try: pic(s, ASSET+"logo_h_white.png", 4.92, 0.9, w=3.5)
except Exception: pass
text(s,1.0,2.4,11.33,0.9,[P(("MERCI",52,WHITE,True))],align=PP_ALIGN.CENTER)
text(s,1.0,3.45,11.33,0.4,[P(("Azélie Bernard   ·   Sébastien Lazcanotegui",17,BEIGE_TX,False))],align=PP_ALIGN.CENTER)
text(s,1.0,3.85,11.33,0.35,[P(("Master 1 Big Data & Intelligence Artificielle — 2026",13,BLUEGREY,False))],align=PP_ALIGN.CENTER)
qr=[("Repository GitHub"),("Rapport complet"),("Model Card")]
for i,lbl in enumerate(qr):
    x=2.85+i*2.7
    rect(s,x,4.55,1.4,1.4,fill=WHITE,shape=MSO_SHAPE.ROUNDED_RECTANGLE,radius=0.08)
    text(s,x,4.55,1.4,1.4,[P(("QR",20,PETROL,True))],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x-0.3,6.0,2.0,0.35,[P((lbl,12,BEIGE_TX,True))],align=PP_ALIGN.CENTER)
text(s,1.0,6.55,11.33,0.4,[P(("Questions ? Disponible dans le chat de soutenance.",13,BLUEGREY,False,True))],align=PP_ALIGN.CENTER)

prs.save(OUT)
print("SAVED", OUT, os.path.getsize(OUT), "bytes,", len(prs.slides._sldIdLst), "slides")
